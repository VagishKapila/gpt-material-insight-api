# utils/scope_parser.py
import os
import docx
import fitz  # PyMuPDF
import openpyxl
import pptx

def parse_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        return [line.strip() for line in f if line.strip()]

def parse_pdf(file_path):
    doc = fitz.open(file_path)
    lines = []
    for page in doc:
        lines.extend(page.get_text().splitlines())
    return [line.strip() for line in lines if line.strip()]

def parse_docx(file_path):
    doc = docx.Document(file_path)
    return [para.text.strip() for para in doc.paragraphs if para.text.strip()]

def parse_xlsx(file_path):
    wb = openpyxl.load_workbook(file_path)
    sheet = wb.active
    lines = []
    for row in sheet.iter_rows(values_only=True):
        for cell in row:
            if cell and isinstance(cell, str):
                lines.append(cell.strip())
    return lines

def parse_pptx(file_path):
    prs = pptx.Presentation(file_path)
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                lines.append(shape.text.strip())
    return [line for line in lines if line]

def parse_scope_file(path):
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == ".txt":
            return parse_txt(path)
        elif ext == ".pdf":
            return parse_pdf(path)
        elif ext == ".docx":
            return parse_docx(path)
        elif ext == ".xlsx":
            return parse_xlsx(path)
        elif ext == ".pptx":
            return parse_pptx(path)
        else:
            return []
    except Exception as e:
        print(f"❌ Failed to parse {ext}: {e}")
        return []
